#!/usr/bin/env python3
"""
Script for converting binary files to text with OCR, translation, and AI image analysis
Supported formats: DOCX, PDF, XLSX, XLS, PPTX, MSG, PFILE, PNG, EML
Features:
- Automatic language detection and translation to English
- OCR for images
- AI vision analysis for images using Google Gemini (primary service)
- Cost estimation before API calls
"""

import os
import sys
import shutil
import threading
import hashlib
import json
import re
import zipfile
from html import unescape
from pathlib import Path
import logging
from concurrent.futures import ThreadPoolExecutor, TimeoutError, as_completed
from typing import Dict, Any, Optional
import tempfile
import time
import warnings
from app.services.evidence.utils import (
    DEFAULT_MAX_RETRIES,
    DEFAULT_RETRY_BASE_DELAY_S,
    IMAGE_ANALYSIS_TIMEOUT_S,
    MAX_RETRY_DELAY_S,
    configure_utf8_output,
    load_env_file,
)


try:
    from docx import Document
    import openpyxl
    from pptx import Presentation
    import extract_msg
    import xlrd
    from PIL import Image
    import pytesseract
    import email
    from deep_translator import GoogleTranslator
    from langdetect import detect, LangDetectException
    CLAUDE_AVAILABLE = False
    try:
        import anthropic as anthropic_lib
        CLAUDE_AVAILABLE = True
    except ImportError:
        pass  # Claude will be unavailable
except ImportError as e:
    print(f"Error importing libraries: {e}")
    print("Please install missing packages: pip install deep-translator langdetect anthropic")
    sys.exit(1)

logger = logging.getLogger(__name__)

# API Configuration — loaded from .env / environment, NOT hardcoded
CLAUDE_API_KEY = os.environ.get("CLAUDE_API_KEY", "")

# Cost constants (as of 2025)
CLAUDE_VISION_COST_PER_1K_TOKENS = 0.00025  # $0.00025 per 1K tokens for Claude Haiku (cost-effective)

# Log file extensions that should not be translated
LOG_EXTENSIONS = {'.log', '.actt'}
METADATA_BEGIN = "=== AUDIT_TOOL_SOURCE_METADATA ==="
METADATA_END = "=== /AUDIT_TOOL_SOURCE_METADATA ==="
TRANSLATION_MAX_CHARS = 4500
TRANSLATION_BLOCK_SEPARATOR = "\n[[[AUDIT_TOOL_BLOCK_SPLIT]]]\n"
TRANSLATION_RETRY_BASE_DELAY_S = 2.0
TRANSLATION_MAX_RETRIES = 3
TRANSLATION_RETRY_MAX_DELAY_S = 8.0
TRANSLATION_FAILURE_COOLDOWN_S = 300.0
CONVERSION_STATE_FILE = "_conversion_state.json"
MARKUP_EXTENSIONS = {'.html', '.htm', '.drawio'}
ZIP_MAX_MEMBERS = 500
ZIP_MAX_TOTAL_BYTES = 250 * 1024 * 1024
MAX_EMBEDDED_RECURSION = 2


def safe_attachment_name(name: Optional[str], fallback_prefix: str = "attachment") -> str:
    """Return a filesystem-safe attachment filename that stays within temp dirs."""
    raw_name = (name or "").strip()
    candidate = Path(raw_name).name
    if not candidate or candidate in {".", ".."}:
        candidate = fallback_prefix
    safe = "".join(ch if ch.isalnum() or ch in "._- " else "_" for ch in candidate).strip(" .")
    return safe or fallback_prefix


class TranslationService:
    """Service for language detection and translation using Google Translate - Always translates to English only"""
    
    def __init__(self, warning_callback=None):
        # Always translate to English only - target='en' is hardcoded
        self.translator = GoogleTranslator(source='auto', target='en')
        self.translation_cache = {}  # Cache translations to avoid duplicate API calls
        self._stats_lock = threading.Lock()
        self._warning_callback = warning_callback
        self._unavailable_until = 0.0
        self._unavailable_notice_emitted = False
        self.stats = {
            'total_text_checked': 0,
            'text_translated': 0,
            'text_already_english': 0,
            'translation_errors': 0,
            'translation_retries': 0,
            'translation_skipped_unavailable': 0,
        }

    @staticmethod
    def _cache_key(text: str, file_extension: str) -> str:
        normalized = text.strip().replace('\r\n', '\n')
        digest = hashlib.sha256()
        digest.update(file_extension.lower().encode('utf-8'))
        digest.update(b'\0')
        digest.update(normalized.encode('utf-8', errors='ignore'))
        return digest.hexdigest()

    @staticmethod
    def _contains_non_ascii_letters(text: str) -> bool:
        return any(ch.isalpha() and ord(ch) > 127 for ch in text)

    def _should_skip_short_translation(self, text: str) -> bool:
        stripped = text.strip()
        return not any(ch.isalpha() for ch in stripped)

    def _mark_warning(self):
        if self._warning_callback:
            self._warning_callback()

    def _translation_is_unavailable(self) -> bool:
        with self._stats_lock:
            if time.monotonic() >= self._unavailable_until:
                self._unavailable_notice_emitted = False
                return False
            self.stats['translation_skipped_unavailable'] += 1
            should_log = not self._unavailable_notice_emitted
            self._unavailable_notice_emitted = True
        self._mark_warning()
        if should_log:
            logger.warning(
                "Translation service is temporarily unavailable; preserving original text "
                "for additional content in this run."
            )
        return True

    def _suspend_translation(self, exc: Exception):
        with self._stats_lock:
            self._unavailable_until = time.monotonic() + TRANSLATION_FAILURE_COOLDOWN_S
            self._unavailable_notice_emitted = False
        logger.warning(
            "Translation service unavailable after retries; preserving original text "
            "and pausing translation attempts for %.0fs: %s",
            TRANSLATION_FAILURE_COOLDOWN_S,
            exc,
        )

    def _split_text_into_chunks(self, text: str, max_chunk_size: int = TRANSLATION_MAX_CHARS):
        chunks = []
        current = ""
        for line in text.splitlines(keepends=True):
            if len(line) > max_chunk_size:
                if current:
                    chunks.append(current)
                    current = ""
                for idx in range(0, len(line), max_chunk_size):
                    chunks.append(line[idx:idx + max_chunk_size])
                continue
            if current and len(current) + len(line) > max_chunk_size:
                chunks.append(current)
                current = line
            else:
                current += line
        if current:
            chunks.append(current)
        return chunks or [text]

    @staticmethod
    def _is_transient_translation_error(err_str: str) -> bool:
        lowered = err_str.lower()
        transient_markers = (
            "request exception",
            "connection error",
            "timed out",
            "timeout",
            "temporarily unavailable",
            "503",
            "429",
        )
        return any(marker in lowered for marker in transient_markers)

    def _translate_chunk_with_retry(self, chunk: str) -> str:
        delay = TRANSLATION_RETRY_BASE_DELAY_S
        last_err = None
        for attempt in range(TRANSLATION_MAX_RETRIES):
            try:
                return self.translator.translate(chunk)
            except Exception as exc:
                last_err = exc
                if self._is_transient_translation_error(str(exc)) and attempt < TRANSLATION_MAX_RETRIES - 1:
                    with self._stats_lock:
                        self.stats['translation_retries'] += 1
                    logger.warning(
                        "Translation transient error, retrying in %.0fs (attempt %s/%s): %s",
                        delay,
                        attempt + 1,
                        TRANSLATION_MAX_RETRIES,
                        exc,
                    )
                    time.sleep(delay)
                    delay = min(delay * 2, TRANSLATION_RETRY_MAX_DELAY_S)
                    continue
                raise last_err
        raise last_err
    
    def detect_language(self, text: str) -> Optional[str]:
        """Detect the language of the text"""
        if not text or len(text.strip()) < 3:
            return None
        try:
            # Sample first 1000 chars for detection (langdetect works better with more text)
            sample = text[:1000] if len(text) > 1000 else text
            detected = detect(sample)
            return detected
        except LangDetectException:
            return None
        except Exception as e:
            logger.warning(f"Language detection error: {e}")
            return None
    
    def is_english(self, text: str) -> bool:
        """Check if text is already in English"""
        lang = self.detect_language(text)
        return lang == 'en' if lang else not self._contains_non_ascii_letters(text)
    
    def translate_to_english(self, text: str, file_extension: str = '') -> str:
        """
        Translate text to English if it's not already in English.
        IMPORTANT: This method ALWAYS translates to English only (target='en' is hardcoded in __init__).
        """
        if not text or len(text.strip()) < 3:
            return text
        
        # Skip translation for log files
        if file_extension.lower() in LOG_EXTENSIONS:
            return text
        
        with self._stats_lock:
            self.stats['total_text_checked'] += 1
        
        # Check cache first
        cache_key = self._cache_key(text, file_extension)
        cached = self.translation_cache.get(cache_key)
        if cached is not None:
            return cached

        if self._should_skip_short_translation(text):
            with self._stats_lock:
                self.stats['text_already_english'] += 1
            self.translation_cache[cache_key] = text
            return text
        
        # Check if already English
        if self.is_english(text):
            with self._stats_lock:
                self.stats['text_already_english'] += 1
            self.translation_cache[cache_key] = text
            return text

        if self._translation_is_unavailable():
            return text
        
        # Translate to English only
        # Note: self.translator is initialized with target='en' in __init__,
        # so all translate() calls will automatically translate to English
        had_error = False
        try:
            translated_chunks = []
            chunks = self._split_text_into_chunks(text)
            for index, chunk in enumerate(chunks):
                if not chunk.strip():
                    translated_chunks.append(chunk)
                    continue
                try:
                    translated_chunks.append(self._translate_chunk_with_retry(chunk))
                except Exception as exc:
                    logger.error(f"Translation error: {exc}")
                    translated_chunks.append(chunk)
                    translated_chunks.extend(chunks[index + 1:])
                    if self._is_transient_translation_error(str(exc)):
                        self._suspend_translation(exc)
                    had_error = True
                    break
            result = "".join(translated_chunks)

            self.translation_cache[cache_key] = result
            with self._stats_lock:
                if had_error:
                    self.stats['translation_errors'] += 1
                    self._mark_warning()
                else:
                    self.stats['text_translated'] += 1
            logger.info(f"Translated text (length: {len(text)}) to English")
            return result
        except Exception as e:
            logger.error(f"Translation error: {e}")
            with self._stats_lock:
                self.stats['translation_errors'] += 1
            self._mark_warning()
            return text  # Return original on error

    def translate_blocks(self, blocks, file_extension: str = ''):
        """Translate ordered document blocks in larger bundles to reduce API overhead."""
        translated_blocks = []
        bundle = []
        bundle_chars = 0

        def flush_bundle():
            nonlocal bundle, bundle_chars
            if not bundle:
                return
            if len(bundle) == 1:
                translated_blocks.append(self.translate_to_english(bundle[0], file_extension))
            else:
                merged = TRANSLATION_BLOCK_SEPARATOR.join(bundle)
                translated = self.translate_to_english(merged, file_extension)
                parts = translated.split(TRANSLATION_BLOCK_SEPARATOR)
                if len(parts) == len(bundle):
                    translated_blocks.extend(parts)
                else:
                    translated_blocks.extend(
                        self.translate_to_english(block, file_extension) for block in bundle
                    )
            bundle = []
            bundle_chars = 0

        for block in blocks:
            if not block or not block.strip():
                flush_bundle()
                translated_blocks.append(block)
                continue
            if len(block) > TRANSLATION_MAX_CHARS:
                flush_bundle()
                translated_blocks.append(self.translate_to_english(block, file_extension))
                continue
            projected = bundle_chars + len(block)
            if bundle and projected + len(TRANSLATION_BLOCK_SEPARATOR) > TRANSLATION_MAX_CHARS:
                flush_bundle()
            bundle.append(block)
            bundle_chars += len(block)
        flush_bundle()
        return translated_blocks
    
    def get_stats(self) -> Dict[str, int]:
        """Get translation statistics"""
        with self._stats_lock:
            return self.stats.copy()


class AIVisionService:
    """Service for AI image analysis using Anthropic Claude (anthropic SDK).

    No fixed inter-request delay — back-to-back calls are fine.
    On a genuine 429 we back off and retry with exponential delay.
    """

    # claude-haiku-4-5: cheaper/faster for bulk image extraction workloads.
    # categorize_by_subject.py uses claude-haiku for text reasoning too.
    MODEL_NAME: str = "claude-haiku-4-5-20251001"
    RETRY_BASE_DELAY_S: float = DEFAULT_RETRY_BASE_DELAY_S
    MAX_RETRIES: int = DEFAULT_MAX_RETRIES

    def __init__(self, warning_callback=None):
        self.client = None
        self.available = False
        self.provider = None
        self._warning_callback = warning_callback
        self._stats_lock = threading.Lock()  # thread-safe stats for parallel processing

        if CLAUDE_AVAILABLE and CLAUDE_API_KEY:
            try:
                import anthropic as anthropic_lib
                self.client = anthropic_lib.Anthropic(api_key=CLAUDE_API_KEY)
                self.provider = 'claude'
                self.available = True
                logger.info(f"Claude client ready (model: {self.MODEL_NAME})")
            except Exception as e:
                logger.warning(f"Claude init failed: {e}. AI vision disabled.")

        if not self.available:
            logger.info("Claude AI vision not available — OCR fallback will be used.")

        self.stats = {
            'images_processed': 0,
            'api_errors': 0,
            'total_cost_estimated': 0.0,
            'total_cost_actual': 0.0,
            'provider_used': self.provider,
        }

    def estimate_image_cost(self, image_path: Path) -> float:
        try:
            img = Image.open(image_path)
            pixels = img.size[0] * img.size[1]
            if pixels < 512 * 512:
                return 0.00005
            elif pixels < 1024 * 1024:
                return 0.0001
            return 0.00025
        except Exception:
            return 0.0001

    def analyze_image(self, image_path: Path, ocr_text: str = "") -> Optional[str]:
        if not self.available or self.client is None:
            return None

        prompt = (
            "Analyze this image and extract ALL information. Translate everything to English.\n"
            "1. Extract all visible text verbatim\n"
            "2. Preserve tables with their structure and numbers\n"
            "3. Describe charts/graphs with their data values\n"
            "4. Note any diagrams, labels, or UI elements\n"
            "Be thorough — this is for an audit report."
        )
        if ocr_text:
            prompt += f"\nOCR hint (may be partial): {ocr_text[:500]}"

        delay = self.RETRY_BASE_DELAY_S
        last_err = None

        for attempt in range(self.MAX_RETRIES):
            try:
                result = self._call_claude(image_path, prompt)
                return result
            except Exception as e:
                last_err = e
                err_str = str(e)
                lowered = err_str.lower()
                is_transient = (
                    "429" in err_str
                    or "503" in err_str
                    or "quota" in lowered
                    or "overloaded" in lowered
                    or "unavailable" in lowered
                    or "timeout" in lowered
                    or "timed out" in lowered
                )
                if is_transient and attempt < self.MAX_RETRIES - 1:
                    logger.warning(
                        f"Transient Claude error for {image_path.name} — "
                        f"waiting {delay:.0f}s (attempt {attempt+1}/{self.MAX_RETRIES})"
                    )
                    time.sleep(delay)
                    delay = min(delay * 2, MAX_RETRY_DELAY_S)
                else:
                    logger.error(f"Claude error for {image_path.name}: {e}")
                    break

        logger.error(f"Claude vision failed for {image_path.name}: {last_err}")
        with self._stats_lock:
            self.stats['api_errors'] += 1
        if self._warning_callback:
            self._warning_callback()
        return None

    def _call_claude(self, image_path: Path, prompt: str) -> str:
        import base64
        import io
        import PIL.Image as PILImage

        with self._stats_lock:
            self.stats['total_cost_estimated'] += 0.0001

        img = PILImage.open(image_path)
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        img_b64 = base64.standard_b64encode(buf.getvalue()).decode("ascii")

        with ThreadPoolExecutor(max_workers=1) as pool:
            future = pool.submit(
                self.client.messages.create,
                model=self.MODEL_NAME,
                max_tokens=1024,
                messages=[{
                    "role": "user",
                    "content": [
                        {"type": "image", "source": {"type": "base64", "media_type": "image/png", "data": img_b64}},
                        {"type": "text", "text": prompt},
                    ],
                }],
            )
            try:
                response = future.result(timeout=IMAGE_ANALYSIS_TIMEOUT_S)
            except TimeoutError as exc:
                future.cancel()
                raise RuntimeError(
                    f"Claude image analysis timed out after {IMAGE_ANALYSIS_TIMEOUT_S:.0f}s"
                ) from exc

        text = response.content[0].text
        with self._stats_lock:
            self.stats['images_processed'] += 1
            try:
                tokens = response.usage.input_tokens + response.usage.output_tokens
                self.stats['total_cost_actual'] += (tokens / 1000) * CLAUDE_VISION_COST_PER_1K_TOKENS
            except Exception:
                pass
        logger.info(f"Claude analyzed: {image_path.name}")
        return text

    def get_stats(self) -> Dict[str, Any]:
        with self._stats_lock:
            return self.stats.copy()


class OCRService:
    """OCR wrapper with one-time capability detection and stats."""

    def __init__(self, warning_callback=None):
        self.available = shutil.which("tesseract") is not None
        self._warned_missing = False
        self._stats_lock = threading.Lock()
        self._warning_callback = warning_callback
        self.stats = {
            'ocr_attempts': 0,
            'ocr_errors': 0,
            'ocr_unavailable': 0,
        }
        if not self.available:
            logger.warning("Tesseract OCR not found on PATH. OCR fallback is disabled.")

    def extract_text(self, image_path: Path) -> str:
        with self._stats_lock:
            self.stats['ocr_attempts'] += 1

        if not self.available:
            with self._stats_lock:
                self.stats['ocr_unavailable'] += 1
            if self._warning_callback:
                self._warning_callback()
            if not self._warned_missing:
                logger.warning("Tesseract OCR is not installed or not on PATH. OCR fallback will be skipped.")
                self._warned_missing = True
            return ""

        try:
            img = Image.open(image_path)
            text = pytesseract.image_to_string(img, lang='eng+rus+chi_sim+spa+heb')
            return text.strip()
        except Exception as e:
            with self._stats_lock:
                self.stats['ocr_errors'] += 1
            if self._warning_callback:
                self._warning_callback()
            logger.error(f"OCR error for {image_path}: {e}")
            return ""

    def get_stats(self) -> Dict[str, int]:
        with self._stats_lock:
            return self.stats.copy()


def save_and_process_image(img_bytes: bytes, img_ext: str, ocr_prefix: str, 
                           temp_dir: str, translation_service: TranslationService,
                           ocr_service: OCRService,
                           vision_service: AIVisionService) -> str:
    """Save image, try AI Vision Analysis first, fallback to OCR if AI Vision fails or unavailable"""
    image_path = None
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=img_ext, dir=temp_dir) as tmp_img:
            tmp_img.write(img_bytes)
            tmp_img.flush()
            image_path = Path(tmp_img.name)
        
        # Step 1: Try AI Vision Analysis first (primary method)
        ai_analysis = None
        if vision_service.available:
            ai_analysis = vision_service.analyze_image(image_path, "")
        
        # Step 2: Only use OCR as fallback if AI Vision failed or is unavailable
        ocr_text = ""
        translated_ocr = ""
        if not ai_analysis:
            # AI Vision failed or unavailable, fallback to OCR
            if vision_service.available:
                logger.info(f"AI Vision failed for {ocr_prefix}, falling back to OCR")
            else:
                logger.info(f"AI Vision unavailable for {ocr_prefix}, using OCR")
            ocr_text = ocr_service.extract_text(image_path)
            if ocr_text:
                translated_ocr = translation_service.translate_to_english(ocr_text)
        
        # Combine results - prioritize AI Vision Analysis
        result_parts = []
        if ai_analysis:
            result_parts.append(f"\n[{ocr_prefix} - AI Vision Analysis]\n{ai_analysis}\n")
        elif translated_ocr:
            # Only show OCR if AI Vision didn't work
            result_parts.append(f"\n[{ocr_prefix} - OCR Text (Translated)]\n{translated_ocr}\n")
        
        return "\n".join(result_parts) if result_parts else ""
        
    except Exception as e:
        logger.error(f"Error processing image: {e}")
        return ""
    finally:
        if image_path and image_path.exists():
            image_path.unlink(missing_ok=True)


class FileConverter:
    PLAIN_TEXT_EXTENSIONS = {
        '.txt', '.actt', '.log', '.ps1', '.csv', '.json', '.xml', '.cfg',
        '.ini', '.conf', '.yml', '.yaml', '.md', '.html', '.htm', '.drawio'
    }

    def __init__(self, input_dir: str, output_dir: str):
        self.input_dir = Path(input_dir)
        self.output_dir = Path(output_dir)
        self._warning_context = threading.local()
        self.translation_service = TranslationService(self._mark_current_file_warning)
        self.vision_service = AIVisionService(self._mark_current_file_warning)
        self.ocr_service = OCRService(self._mark_current_file_warning)
        self.total_estimated_cost = 0.0
        self._cost_lock = threading.Lock()
        self.state_file = self.output_dir / CONVERSION_STATE_FILE
        self._state_lock = threading.Lock()
        self._conversion_state = self._load_conversion_state()
        self._state_dirty = False

    def _add_estimated_cost(self, cost: float):
        with self._cost_lock:
            self.total_estimated_cost += cost

    def _estimate_embedded_image_cost(self, img_bytes: bytes, img_ext: str, temp_dir: str) -> float:
        image_path = None
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix=img_ext, dir=temp_dir) as tmp:
                tmp.write(img_bytes)
                tmp.flush()
                image_path = Path(tmp.name)
            return self.vision_service.estimate_image_cost(image_path)
        finally:
            if image_path and image_path.exists():
                image_path.unlink(missing_ok=True)

    def _output_path_for(self, file_path: Path) -> Path:
        relative_path = file_path.relative_to(self.input_dir)
        return self.output_dir / relative_path.with_suffix('.txt')

    def _source_sha256(self, file_path: Path) -> str:
        digest = hashlib.sha256()
        with open(file_path, 'rb') as handle:
            for chunk in iter(lambda: handle.read(1024 * 1024), b''):
                digest.update(chunk)
        return digest.hexdigest()

    def _source_metadata(self, file_path: Path) -> Dict[str, str]:
        stat = file_path.stat()
        return {
            'source_path': str(file_path.relative_to(self.input_dir)),
            'source_name': file_path.name,
            'source_size': str(stat.st_size),
            'source_mtime_ns': str(stat.st_mtime_ns),
            'source_sha256': self._source_sha256(file_path),
        }

    def _load_conversion_state(self) -> Dict[str, Dict[str, str]]:
        if not self.state_file.exists():
            return {}
        try:
            raw = json.loads(self.state_file.read_text(encoding='utf-8'))
            files = raw.get('files', {})
            return files if isinstance(files, dict) else {}
        except Exception as e:
            logger.warning(f"Could not read conversion state from {self.state_file}: {e}")
            return {}

    def _save_conversion_state(self):
        with self._state_lock:
            if not self._state_dirty:
                return
            payload = {
                'version': 1,
                'files': self._conversion_state,
            }
            self.state_file.parent.mkdir(parents=True, exist_ok=True)
            self.state_file.write_text(
                json.dumps(payload, indent=2, ensure_ascii=False),
                encoding='utf-8',
            )
            self._state_dirty = False

    def _update_conversion_state(self, file_path: Path, processing_complete: bool = True):
        key = str(file_path.relative_to(self.input_dir))
        metadata = self._source_metadata(file_path).copy()
        metadata['output_path'] = str(self._output_path_for(file_path).relative_to(self.output_dir))
        metadata['processing_complete'] = processing_complete
        with self._state_lock:
            self._conversion_state[key] = metadata
            self._state_dirty = True

    def _format_output_payload(self, file_path: Path, text: str, processing_complete: bool = True) -> str:
        metadata = self._source_metadata(file_path)
        if not processing_complete:
            metadata['processing_complete'] = 'false'
        header_lines = [METADATA_BEGIN]
        for key, value in metadata.items():
            header_lines.append(f"{key}: {value}")
        header_lines += [METADATA_END, ""]
        return "\n".join(header_lines) + text

    def _read_output_metadata(self, output_file: Path) -> Optional[Dict[str, str]]:
        if not output_file.exists():
            return None
        try:
            with open(output_file, 'r', encoding='utf-8', errors='ignore') as handle:
                lines = []
                for _ in range(16):
                    line = handle.readline()
                    if not line:
                        break
                    lines.append(line.rstrip('\n'))
                    if line.rstrip('\n') == METADATA_END:
                        break
        except Exception:
            return None

        if not lines or lines[0] != METADATA_BEGIN or METADATA_END not in lines:
            return None

        metadata: Dict[str, str] = {}
        for line in lines[1:]:
            if line == METADATA_END:
                break
            if ": " not in line:
                continue
            key, value = line.split(": ", 1)
            metadata[key] = value
        return metadata

    def _is_up_to_date(self, file_path: Path) -> bool:
        output_file = self._output_path_for(file_path)
        current = self._source_metadata(file_path)
        state_key = str(file_path.relative_to(self.input_dir))
        state_entry = self._conversion_state.get(state_key)
        if state_entry and output_file.exists():
            if state_entry.get('processing_complete') is False:
                return False
            return all(state_entry.get(key) == value for key, value in current.items())
        existing = self._read_output_metadata(output_file)
        if existing is None:
            return False
        if existing.get('processing_complete', '').lower() == 'false':
            return False
        return all(existing.get(key) == value for key, value in current.items())

    @staticmethod
    def _is_probably_text_file(file_path: Path) -> bool:
        try:
            with open(file_path, 'rb') as handle:
                sample = handle.read(4096)
        except Exception:
            return False
        if not sample:
            return True
        if b'\x00' in sample:
            return False
        try:
            sample.decode('utf-8')
            return True
        except UnicodeDecodeError:
            try:
                sample.decode('cp1252')
                return True
            except UnicodeDecodeError:
                return False

    def estimate_processing_costs(self, file_path: Path) -> float:
        """Estimate total costs for processing a file"""
        estimated_cost = 0.0
        extension = file_path.suffix.lower()
        
        # Estimate image processing costs
        if extension in ['.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.gif']:
            estimated_cost += self.vision_service.estimate_image_cost(file_path)
        elif extension in ['.docx', '.pdf', '.pptx', '.msg', '.eml']:
            # These may contain images - rough estimate
            # We'll refine this during actual processing
            pass
        
        return estimated_cost

    def translate_text_content(self, text: str, file_extension: str = '') -> str:
        """Translate text content to English if needed"""
        if not text:
            return text
        return self.translation_service.translate_to_english(text, file_extension)

    def _begin_file_warning_tracking(self):
        self._warning_context.had_warning = False

    def _mark_current_file_warning(self):
        self._warning_context.had_warning = True

    def _current_file_had_warning(self) -> bool:
        return bool(getattr(self._warning_context, 'had_warning', False))

    def _format_table(self, table) -> str:
        """Format a DOCX table as readable text"""
        table_text = []
        for row in table.rows:
            row_cells = []
            for cell in row.cells:
                cell_text = cell.text.strip().replace('\n', ' ').replace('\r', ' ')
                row_cells.append(cell_text)
            # Join cells with tab separator for readability
            table_text.append('\t'.join(row_cells))
        return '\n'.join(table_text)

    @staticmethod
    def _extract_markup_text(content: str) -> str:
        content = re.sub(r'(?is)<(script|style).*?>.*?</\1>', ' ', content)
        content = re.sub(r'(?i)<br\s*/?>', '\n', content)
        content = re.sub(r'(?i)</(p|div|tr|li|h[1-6])>', '\n', content)
        content = re.sub(r'<[^>]+>', ' ', content)
        content = unescape(content)
        content = re.sub(r'\r\n?', '\n', content)
        content = re.sub(r'[ \t]+', ' ', content)
        content = re.sub(r'\n{3,}', '\n\n', content)
        return content.strip()

    def _read_plain_text_for_conversion(self, file_path: Path) -> str:
        with open(file_path, 'r', encoding='utf-8-sig', errors='ignore') as f:
            content = f.read()
        if file_path.suffix.lower() in MARKUP_EXTENSIONS:
            content = self._extract_markup_text(content)
        return content

    def _convert_embedded_file_to_text(self, file_path: Path, temp_dir: str, recursion_depth: int = 0) -> Optional[str]:
        suffix = file_path.suffix.lower()
        if suffix in self.PLAIN_TEXT_EXTENSIONS or (not suffix and self._is_probably_text_file(file_path)):
            content = self._read_plain_text_for_conversion(file_path)
            if suffix not in LOG_EXTENSIONS:
                content = self.translate_text_content(content, suffix)
            return content
        return self.convert_file(file_path, temp_dir, recursion_depth=recursion_depth)

    def convert_zip(self, file_path: Path, temp_dir: str, recursion_depth: int = 0) -> Optional[str]:
        """Convert ZIP archive by extracting supported member files recursively."""
        try:
            if recursion_depth >= MAX_EMBEDDED_RECURSION:
                return f"=== ZIP ARCHIVE: {file_path.name} ===\n[ZIP processing skipped: maximum embedded recursion depth reached]"
            extract_root = (Path(temp_dir) / f"zip_{file_path.stem}").resolve()
            extract_root.mkdir(parents=True, exist_ok=True)
            text_parts = [f"=== ZIP ARCHIVE: {file_path.name} ==="]
            total_bytes = 0
            member_count = 0

            with zipfile.ZipFile(file_path) as archive:
                for member in archive.infolist():
                    if member.is_dir():
                        continue
                    member_count += 1
                    total_bytes += member.file_size
                    if member_count > ZIP_MAX_MEMBERS:
                        text_parts.append("[ZIP processing stopped: too many files in archive]")
                        break
                    if total_bytes > ZIP_MAX_TOTAL_BYTES:
                        text_parts.append("[ZIP processing stopped: archive exceeds size safety limit]")
                        break

                    member_path = Path(member.filename)
                    if member_path.is_absolute() or member_path.drive or ".." in member_path.parts:
                        text_parts.append(f"[ZIP member skipped: unsafe path {member.filename}]")
                        continue
                    safe_parts = [part for part in member_path.parts if part not in ("", ".")]
                    if not safe_parts:
                        continue
                    target_path = extract_root.joinpath(*safe_parts).resolve()
                    try:
                        target_path.relative_to(extract_root)
                    except ValueError:
                        text_parts.append(f"[ZIP member skipped: unsafe path {member.filename}]")
                        continue
                    target_path.parent.mkdir(parents=True, exist_ok=True)
                    with archive.open(member) as source, open(target_path, 'wb') as dest:
                        shutil.copyfileobj(source, dest)

                    suffix = target_path.suffix.lower()
                    if suffix not in self.PLAIN_TEXT_EXTENSIONS and suffix not in {'.docx', '.pdf', '.xlsx', '.xls', '.pptx', '.msg', '.pfile', '.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.gif', '.eml', '.zip'}:
                        if not (not suffix and self._is_probably_text_file(target_path)):
                            continue

                    converted = self._convert_embedded_file_to_text(target_path, temp_dir, recursion_depth=recursion_depth + 1)
                    if converted:
                        text_parts.extend([
                            "",
                            f"--- ZIP MEMBER: {member.filename} ---",
                            converted,
                        ])
            return "\n".join(text_parts)
        except Exception as e:
            logger.error(f"Error converting ZIP {file_path}: {e}")
            return None
    
    def convert_docx(self, file_path: Path, temp_dir: str) -> Optional[str]:
        """Convert DOCX file with translation and AI image analysis"""
        try:
            doc = Document(file_path)
            body_blocks = []
            external_links = []
            
            # Extract content in order: paragraphs and tables
            # Create a mapping of element to table/paragraph for quick lookup
            table_map = {table._element: table for table in doc.tables}
            paragraph_map = {p._element: p for p in doc.paragraphs}
            
            # Iterate through document body elements to maintain order
            for element in doc.element.body:
                # Check if element is a paragraph (handle different namespace formats)
                tag_local = element.tag.split('}')[-1] if '}' in element.tag else element.tag
                if tag_local == 'p':
                    paragraph = paragraph_map.get(element)
                    if paragraph and paragraph.text.strip():
                        body_blocks.append(paragraph.text)
                # Check if element is a table
                elif tag_local == 'tbl':
                    table = table_map.get(element)
                    if table:
                        formatted_table = self._format_table(table)
                        if formatted_table.strip():
                            body_blocks.append(formatted_table)

            text = self.translation_service.translate_blocks(body_blocks, '.docx')
            
            # Extract and process images from all document parts
            ocr_results = []
            rel_idx = 0
            processed_image_ids = set()  # Track processed images to avoid duplicates
            
            # Process images from main document part
            for rel in doc.part.rels.values():
                if getattr(rel, 'target_mode', None) == 'External':
                    external_links.append(rel.target_ref)
                elif "image" in rel.target_ref and hasattr(rel, 'target_part') and rel.target_part is not None:
                    rel_id = rel.rId
                    if rel_id not in processed_image_ids:
                        processed_image_ids.add(rel_id)
                        rel_idx += 1
                        img_bytes = rel.target_part.blob
                        img_ext = os.path.splitext(rel.target_ref)[-1] or ".png"
                        
                        # Estimate cost before processing
                        cost = self._estimate_embedded_image_cost(img_bytes, img_ext, temp_dir)
                        self._add_estimated_cost(cost)
                        logger.info(f"Estimated cost for DOCX image {rel_idx}: ${cost:.4f}")
                        
                        processed = save_and_process_image(
                            img_bytes, img_ext, f"DOCX image {rel_idx}", temp_dir,
                            self.translation_service, self.ocr_service, self.vision_service
                        )
                        if processed:
                            ocr_results.append(processed)
            
            # Process images from headers and footers
            for section in doc.sections:
                # Header images
                if section.header:
                    for rel in section.header.part.rels.values():
                        if "image" in rel.target_ref and hasattr(rel, 'target_part') and rel.target_part is not None:
                            rel_id = rel.rId
                            if rel_id not in processed_image_ids:
                                processed_image_ids.add(rel_id)
                                rel_idx += 1
                                img_bytes = rel.target_part.blob
                                img_ext = os.path.splitext(rel.target_ref)[-1] or ".png"
                                
                                cost = self._estimate_embedded_image_cost(img_bytes, img_ext, temp_dir)
                                self._add_estimated_cost(cost)
                                logger.info(f"Estimated cost for DOCX header image {rel_idx}: ${cost:.4f}")
                                
                                processed = save_and_process_image(
                                    img_bytes, img_ext, f"DOCX header image {rel_idx}", temp_dir,
                                    self.translation_service, self.ocr_service, self.vision_service
                                )
                                if processed:
                                    ocr_results.append(processed)
                
                # Footer images
                if section.footer:
                    for rel in section.footer.part.rels.values():
                        if "image" in rel.target_ref and hasattr(rel, 'target_part') and rel.target_part is not None:
                            rel_id = rel.rId
                            if rel_id not in processed_image_ids:
                                processed_image_ids.add(rel_id)
                                rel_idx += 1
                                img_bytes = rel.target_part.blob
                                img_ext = os.path.splitext(rel.target_ref)[-1] or ".png"
                                
                                cost = self._estimate_embedded_image_cost(img_bytes, img_ext, temp_dir)
                                self._add_estimated_cost(cost)
                                logger.info(f"Estimated cost for DOCX footer image {rel_idx}: ${cost:.4f}")
                                
                                processed = save_and_process_image(
                                    img_bytes, img_ext, f"DOCX footer image {rel_idx}", temp_dir,
                                    self.translation_service, self.ocr_service, self.vision_service
                                )
                                if processed:
                                    ocr_results.append(processed)
            
            if ocr_results:
                text.append("\n=== IMAGES FROM DOCX ===\n" + "\n".join(ocr_results))
            if external_links:
                text.append("\n=== EXTERNAL LINKS IN DOCX ===\n" + "\n".join(external_links))
                try:
                    ext_log = self.output_dir / "external_links.log"
                    with open(ext_log, "a", encoding="utf-8") as elog:
                        elog.write(f"{file_path}:\n")
                        for link in external_links:
                            elog.write(f"  {link}\n")
                except Exception:
                    pass
            
            # Content already translated per-paragraph and per-table-cell above
            return '\n'.join(text)
        except Exception as e:
            logger.error(f"Error converting DOCX {file_path}: {e}")
            return None

    def convert_pdf(self, file_path: Path, temp_dir: str) -> Optional[str]:
        """Convert PDF file with translation and AI image analysis"""
        try:
            import fitz  # PyMuPDF
        except ImportError:
            logger.error("PyMuPDF (fitz) required for PDF OCR. Install via: pip install pymupdf")
            return None
        try:
            text = []
            ocr_results = []
            doc = fitz.open(str(file_path))
            
            for page_num in range(len(doc)):
                page = doc[page_num]
                page_text = page.get_text()
                if page_text.strip():
                    translated = self.translate_text_content(page_text, '.pdf')
                    text.append(translated)
                
                # Extract and process images
                images = page.get_images(full=True)
                for img_idx, img in enumerate(images, 1):
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    img_bytes = base_image["image"]
                    img_ext = "." + base_image["ext"]
                    
                    # Estimate cost
                    cost = self._estimate_embedded_image_cost(img_bytes, img_ext, temp_dir)
                    self._add_estimated_cost(cost)
                    logger.info(f"Estimated cost for PDF page {page_num+1} image {img_idx}: ${cost:.4f}")
                    
                    processed = save_and_process_image(
                        img_bytes, img_ext, f"PDF page {page_num+1} image {img_idx}", temp_dir,
                        self.translation_service, self.ocr_service, self.vision_service
                    )
                    if processed:
                        ocr_results.append(processed)
            
            if ocr_results:
                text.append("\n=== IMAGES FROM PDF ===\n" + "\n".join(ocr_results))
            
            # Content already translated per-page above
            return '\n'.join(text)
        except Exception as e:
            logger.error(f"Error converting PDF {file_path}: {e}")
            return None

    def convert_pptx(self, file_path: Path, temp_dir: str) -> Optional[str]:
        """Convert PPTX file with translation and AI image analysis"""
        try:
            prs = Presentation(file_path)
            text = []
            ocr_results = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text.append(f"=== Slide {slide_num} ===")
                slide_blocks = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_blocks.append(shape.text)
                    if shape.shape_type == 13:  # PICTURE
                        img = shape.image
                        img_bytes = img.blob
                        img_ext = os.path.splitext(img.ext)[-1] if img.ext else ".png"
                        
                        # Estimate cost
                        cost = self._estimate_embedded_image_cost(img_bytes, img_ext, temp_dir)
                        self._add_estimated_cost(cost)
                        logger.info(f"Estimated cost for PPTX slide {slide_num} image: ${cost:.4f}")
                        
                        processed = save_and_process_image(
                            img_bytes, img_ext, f"PPTX slide {slide_num} image", temp_dir,
                            self.translation_service, self.ocr_service, self.vision_service
                        )
                        if processed:
                            ocr_results.append(processed)
                text.extend(self.translation_service.translate_blocks(slide_blocks, '.pptx'))
                text.append("")
            
            if ocr_results:
                text.append("\n=== IMAGES FROM PPTX ===\n" + "\n".join(ocr_results))
            
            # Content already translated per-shape above
            return '\n'.join(text)
        except Exception as e:
            logger.error(f"Error converting PPTX {file_path}: {e}")
            return None

    def convert_xlsx(self, file_path: Path) -> Optional[str]:
        """Convert XLSX file with translation"""
        try:
            with warnings.catch_warnings(record=True) as captured_warnings:
                warnings.simplefilter("always", UserWarning)
                workbook = openpyxl.load_workbook(file_path, data_only=True)
            for warning in captured_warnings:
                message = str(warning.message)
                if "Data Validation extension is not supported" in message or "Slicer List extension is not supported" in message:
                    logger.info(
                        "Workbook %s contains unsupported Excel UI features; cell values will still be extracted.",
                        file_path.name,
                    )
                    break
            text = []
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text.append(f"=== Sheet: {sheet_name} ===")
                for row in sheet.iter_rows(values_only=True):
                    row_text = [str(cell) for cell in row if cell is not None]
                    if row_text:
                        text.append('\t'.join(row_text))
                text.append("")
            # Translate the entire assembled spreadsheet text in one pass
            result = '\n'.join(text)
            return self.translate_text_content(result, '.xlsx')
        except Exception as e:
            logger.error(f"Error converting XLSX {file_path}: {e}")
            return None

    def convert_xls(self, file_path: Path) -> Optional[str]:
        """Convert XLS file with translation"""
        try:
            workbook = xlrd.open_workbook(file_path)
            text = []
            for sheet_name in workbook.sheet_names():
                sheet = workbook.sheet_by_name(sheet_name)
                text.append(f"=== Sheet: {sheet_name} ===")
                for row_idx in range(sheet.nrows):
                    row_text = [str(sheet.cell_value(row_idx, col_idx))
                               for col_idx in range(sheet.ncols)
                               if sheet.cell_value(row_idx, col_idx)]
                    if row_text:
                        text.append('\t'.join(row_text))
                text.append("")
            # Translate the entire assembled spreadsheet text in one pass
            result = '\n'.join(text)
            return self.translate_text_content(result, '.xls')
        except Exception as e:
            logger.error(f"Error converting XLS {file_path}: {e}")
            return None

    def convert_msg(self, file_path: Path, temp_dir: str, recursion_depth: int = 0) -> Optional[str]:
        """Convert MSG file with translation and AI image analysis"""
        try:
            if recursion_depth >= MAX_EMBEDDED_RECURSION:
                return f"=== MSG MESSAGE: {file_path.name} ===\n[MSG processing skipped: maximum embedded recursion depth reached]"
            msg = extract_msg.Message(file_path)
            text = []
            text.append(f"From: {msg.sender}")
            text.append(f"To: {msg.to}")
            text.append(f"Subject: {self.translate_text_content(msg.subject, '.msg')}")
            text.append(f"Date: {msg.date}")
            text.append("")
            
            if msg.body:
                text.append("=== MESSAGE BODY ===")
                translated_body = self.translate_text_content(msg.body, '.msg')
                text.append(translated_body)
            
            # Process attachments
            if msg.attachments:
                text.append("")
                text.append("=== ATTACHMENTS ===")
                for att in msg.attachments:
                    safe_name = safe_attachment_name(att.longFilename, "msg_attachment")
                    text.append(f"- {safe_name}")
                    att_path = Path(temp_dir) / safe_name
                    with open(att_path, 'wb') as f:
                        f.write(att.data)
                    
                    ext = att_path.suffix.lower()
                    if ext in ['.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.gif']:
                        # Process image with OCR and AI
                        with open(att_path, 'rb') as f:
                            img_bytes = f.read()
                        cost = self.vision_service.estimate_image_cost(att_path)
                        self._add_estimated_cost(cost)
                        logger.info(f"Estimated cost for MSG attachment image: ${cost:.4f}")
                        
                        processed = save_and_process_image(
                            img_bytes, ext, f"MSG attachment {safe_name}", temp_dir,
                            self.translation_service, self.ocr_service, self.vision_service
                        )
                        if processed:
                            text.append(processed)
                    elif self._is_supported_input_file(att_path, {'.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.gif'}, {'.docx', '.pdf', '.xlsx', '.xls', '.pptx', '.msg', '.pfile', '.eml', '.zip'}):
                        conv_text = self._convert_embedded_file_to_text(att_path, temp_dir, recursion_depth=recursion_depth + 1)
                        if conv_text:
                            text.append(f"[Extracted from MSG attachment {safe_name}]\n{conv_text}")
            
            msg.close()
            # Content already translated per-field and body above
            return '\n'.join(text)
        except Exception as e:
            logger.error(f"Error converting MSG {file_path}: {e}")
            return None

    def convert_pfile(self, file_path: Path) -> Optional[str]:
        """Convert PFILE with translation"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            return self.translate_text_content(content, '.pfile')
        except Exception as e:
            logger.error(f"Error converting PFILE {file_path}: {e}")
            return None

    def convert_png(self, file_path: Path, temp_dir: str) -> Optional[str]:
        """Convert PNG/image file with OCR, translation, and AI analysis"""
        try:
            # Estimate cost
            cost = self.vision_service.estimate_image_cost(file_path)
            self._add_estimated_cost(cost)
            logger.info(f"Estimated cost for image {file_path.name}: ${cost:.4f}")
            
            with open(file_path, 'rb') as f:
                img_bytes = f.read()
            
            img_ext = file_path.suffix or ".png"
            processed = save_and_process_image(
                img_bytes, img_ext, f"Image {file_path.name}", temp_dir,
                self.translation_service, self.ocr_service, self.vision_service
            )
            
            if processed:
                return f"=== PROCESSED IMAGE: {file_path.name} ===\n{processed}"
            else:
                return f"=== IMAGE {file_path.name} ===\nFailed to extract information"
        except Exception as e:
            logger.error(f"Error converting PNG {file_path}: {e}")
            return None

    def convert_eml(self, file_path: Path, temp_dir: str, recursion_depth: int = 0) -> Optional[str]:
        """Convert EML file with translation and AI image analysis"""
        try:
            if recursion_depth >= MAX_EMBEDDED_RECURSION:
                return f"=== EML MESSAGE: {file_path.name} ===\n[EML processing skipped: maximum embedded recursion depth reached]"
            with open(file_path, 'rb') as f:
                msg = email.message_from_bytes(f.read())
            
            text = []
            text.append(f"From: {msg.get('From', 'Unknown')}")
            text.append(f"To: {msg.get('To', 'Unknown')}")
            text.append(f"Subject: {self.translate_text_content(msg.get('Subject', 'No subject'), '.eml')}")
            text.append(f"Date: {msg.get('Date', 'Unknown')}")
            text.append("")
            
            if msg.is_multipart():
                text.append("=== MESSAGE BODY ===")
                for part in msg.walk():
                    content_type = part.get_content_type()
                    content_disposition = str(part.get('Content-Disposition', ''))
                    
                    if content_type == 'text/plain' and 'attachment' not in content_disposition:
                        payload = part.get_payload(decode=True)
                        if payload:
                            try:
                                text_content = payload.decode('utf-8', errors='ignore')
                                translated = self.translate_text_content(text_content, '.eml')
                                text.append(translated)
                            except:
                                text.append("Failed to decode text part")
                    elif content_type == 'text/html' and 'attachment' not in content_disposition:
                        payload = part.get_payload(decode=True)
                        if payload:
                            try:
                                html_content = payload.decode('utf-8', errors='ignore')
                                translated = self.translate_text_content(html_content, '.eml')
                                text.append(f"[HTML Content]\n{translated}")
                            except:
                                text.append("Failed to decode HTML part")
                    elif 'attachment' in content_disposition:
                        filename = part.get_filename()
                        if filename:
                            safe_name = safe_attachment_name(filename, "eml_attachment")
                            text.append(f"\n=== ATTACHMENT: {safe_name} ===")
                            att_path = Path(temp_dir) / safe_name
                            try:
                                with open(att_path, 'wb') as att_file:
                                    att_file.write(part.get_payload(decode=True))
                                
                                ext = att_path.suffix.lower()
                                if ext in ['.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.gif']:
                                    with open(att_path, 'rb') as f:
                                        img_bytes = f.read()
                                    cost = self.vision_service.estimate_image_cost(att_path)
                                    self._add_estimated_cost(cost)
                                    logger.info(f"Estimated cost for EML attachment image: ${cost:.4f}")
                                    
                                    processed = save_and_process_image(
                                            img_bytes, ext, f"EML attachment {safe_name}", temp_dir,
                                        self.translation_service, self.ocr_service, self.vision_service
                                    )
                                    if processed:
                                        text.append(processed)
                                elif self._is_supported_input_file(att_path, {'.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.gif'}, {'.docx', '.pdf', '.xlsx', '.xls', '.pptx', '.msg', '.pfile', '.eml', '.zip'}):
                                    conv_text = self._convert_embedded_file_to_text(att_path, temp_dir, recursion_depth=recursion_depth + 1)
                                    if conv_text:
                                        text.append(f"[Extracted from EML attachment {safe_name}]\n{conv_text}")
                            except Exception as e:
                                text.append(f"Error processing attachment {safe_name}: {e}")
            else:
                payload = msg.get_payload(decode=True)
                if payload:
                    try:
                        text_content = payload.decode('utf-8', errors='ignore')
                        translated = self.translate_text_content(text_content, '.eml')
                        text.append("=== MESSAGE BODY ===")
                        text.append(translated)
                    except:
                        text.append("Failed to decode message content")
            
            # Content already translated per-part above
            return '\n'.join(text)
        except Exception as e:
            logger.error(f"Error converting EML {file_path}: {e}")
            return None

    def convert_file(self, file_path: Path, temp_dir: str, recursion_depth: int = 0) -> Optional[str]:
        """Route file to appropriate converter"""
        extension = file_path.suffix.lower()
        if extension == '.docx':
            return self.convert_docx(file_path, temp_dir)
        elif extension == '.pdf':
            return self.convert_pdf(file_path, temp_dir)
        elif extension == '.xlsx':
            return self.convert_xlsx(file_path)
        elif extension == '.xls':
            return self.convert_xls(file_path)
        elif extension == '.pptx':
            return self.convert_pptx(file_path, temp_dir)
        elif extension == '.msg':
            return self.convert_msg(file_path, temp_dir, recursion_depth=recursion_depth)
        elif extension == '.pfile':
            return self.convert_pfile(file_path)
        elif extension in ['.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.gif']:
            return self.convert_png(file_path, temp_dir)
        elif extension == '.eml':
            return self.convert_eml(file_path, temp_dir, recursion_depth=recursion_depth)
        elif extension == '.zip':
            return self.convert_zip(file_path, temp_dir, recursion_depth=recursion_depth)
        else:
            logger.warning(f"Unsupported file format: {file_path}")
            return None

    def copy_plain_file(self, file_path: Path) -> bool:
        """Copy plain-text-like files into normalized .txt outputs."""
        try:
            output_file = self._output_path_for(file_path)
            
            # For non-log text files, translate content
            content = self._read_plain_text_for_conversion(file_path)
            if file_path.suffix.lower() not in LOG_EXTENSIONS:
                content = self.translate_text_content(content, file_path.suffix)
            processing_complete = not self._current_file_had_warning()
            payload = self._format_output_payload(file_path, content, processing_complete)
            output_file.parent.mkdir(parents=True, exist_ok=True)
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(payload)
            self._update_conversion_state(file_path, processing_complete)
            if not processing_complete:
                logger.warning(f"Output retained with warnings and will be retried on the next run: {file_path.name}")
            logger.info(f"Translated and copied text file: {file_path} -> {output_file}")
            return True
        except Exception as e:
            logger.error(f"Error copying text file {file_path}: {e}")
            return False

    def _is_supported_input_file(self, file_path: Path, image_extensions, doc_extensions) -> bool:
        suffix = file_path.suffix.lower()
        if suffix in image_extensions or suffix in doc_extensions or suffix in self.PLAIN_TEXT_EXTENSIONS:
            return True
        if not suffix and self._is_probably_text_file(file_path):
            return True
        return False

    def _iter_input_files(self):
        """Yield all files under input_dir, skipping generated output subtrees."""
        output_resolved = self.output_dir.resolve()
        for file_path in self.input_dir.rglob('*'):
            if not file_path.is_file():
                continue
            relative_parts = file_path.relative_to(self.input_dir).parts
            # Skip any pre-existing generated folders from prior audit runs, even if the
            # current output_dir is configured somewhere else.
            if any(part.lower() == "converted_files" for part in relative_parts):
                continue
            # Also skip files that live inside the actual configured output directory.
            try:
                file_path.resolve().relative_to(output_resolved)
                continue  # it's inside output_dir — skip
            except ValueError:
                pass
            yield file_path

    def convert_all_files(self) -> Dict[str, Any]:
        """Convert all files.

        Stand-alone image files (PNG/JPG/etc.) are processed in parallel via a
        thread pool — Gemini's paid tier handles concurrent requests easily.
        All other formats (DOCX, PDF, XLSX …) are converted sequentially
        because their parsers are not thread-safe.
        """
        # Parallel workers for image-only files.
        # 10 concurrent threads → ~10x throughput on large image batches.
        IMAGE_WORKERS = 10

        stats = {
            'total_files': 0,
            'converted_files': 0,
            'failed_files': 0,
            'skipped_files': 0,
            'unchanged_files': 0,
            'warning_files': 0,
            'unsupported_file_paths': [],
            'failed_file_paths': [],
            'warning_file_paths': [],
        }
        stats_lock = threading.Lock()

        IMAGE_EXTENSIONS = {'.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.gif'}
        DOC_EXTENSIONS   = {'.docx', '.pdf', '.xlsx', '.xls', '.pptx',
                            '.msg', '.pfile', '.eml', '.zip'}
        def _relative_name(file_path: Path) -> str:
            return str(file_path.relative_to(self.input_dir))

        # Partition files into images vs everything else
        image_files  = []
        other_files  = []
        for file_path in self._iter_input_files():
            with stats_lock:
                stats['total_files'] += 1
            suffix = file_path.suffix.lower()
            if not self._is_supported_input_file(file_path, IMAGE_EXTENSIONS, DOC_EXTENSIONS):
                with stats_lock:
                    stats['skipped_files'] += 1
                    stats['unsupported_file_paths'].append(_relative_name(file_path))
                logger.info(f"Skipped (unsupported): {_relative_name(file_path)}")
                continue
            if self._is_up_to_date(file_path):
                with stats_lock:
                    stats['unchanged_files'] += 1
                logger.info(f"Skipped (unchanged): {file_path.name}")
                continue
            if suffix in IMAGE_EXTENSIONS:
                image_files.append(file_path)
            else:
                other_files.append(file_path)

        total_count = len(image_files) + len(other_files)
        logger.info(
            f"Found {total_count} files to process "
            f"({len(image_files)} images [parallel], "
            f"{len(other_files)} documents [sequential])"
        )

        # ── Helper: write one converted text file ───────────────────────
        def _write_output(file_path: Path, text: str, processing_complete: bool):
            output_file = self._output_path_for(file_path)
            output_file.parent.mkdir(parents=True, exist_ok=True)
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(self._format_output_payload(file_path, text, processing_complete))
            self._update_conversion_state(file_path, processing_complete)
            if not processing_complete:
                logger.warning(f"Output retained with warnings and will be retried on the next run: {file_path.name}")
            logger.info(f"Converted: {file_path.name} → {output_file.relative_to(self.output_dir)}")

        # ── 1. Process standalone image files IN PARALLEL ────────────────
        def _convert_single_file(file_path: Path, temp_dir: str):
            self._begin_file_warning_tracking()
            suffix = file_path.suffix.lower()

            if suffix in self.PLAIN_TEXT_EXTENSIONS or (not suffix and self._is_probably_text_file(file_path)):
                ok = self.copy_plain_file(file_path)
            else:
                text = self.convert_file(file_path, temp_dir)
                ok = False
                if text is not None:
                    _write_output(file_path, text, not self._current_file_had_warning())
                    ok = True

            return ok, self._current_file_had_warning()

        if image_files:
            logger.info(f"Processing {len(image_files)} images with {IMAGE_WORKERS} parallel workers…")

            def _process_image(file_path: Path):
                with tempfile.TemporaryDirectory() as td:
                    return _convert_single_file(file_path, td)

            with ThreadPoolExecutor(max_workers=IMAGE_WORKERS) as pool:
                futures = {pool.submit(_process_image, fp): fp for fp in image_files}
                for fut in as_completed(futures):
                    fp = futures[fut]
                    try:
                        ok, had_warnings = fut.result()
                        with stats_lock:
                            if ok:
                                stats['converted_files'] += 1
                                if had_warnings:
                                    stats['warning_files'] += 1
                                    stats['warning_file_paths'].append(_relative_name(fp))
                            else:
                                stats['failed_files'] += 1
                                stats['failed_file_paths'].append(_relative_name(fp))
                    except Exception as e:
                        logger.error(f"Image worker error for {fp.name}: {e}")
                        with stats_lock:
                            stats['failed_files'] += 1
                            stats['failed_file_paths'].append(_relative_name(fp))

        # ── 2. Process document / text files SEQUENTIALLY ────────────────
        with tempfile.TemporaryDirectory() as temp_dir:
            for file_path in other_files:
                ok, had_warnings = _convert_single_file(file_path, temp_dir)
                with stats_lock:
                    if ok:
                        stats['converted_files'] += 1
                        if had_warnings:
                            stats['warning_files'] += 1
                            stats['warning_file_paths'].append(_relative_name(file_path))
                    else:
                        stats['failed_files'] += 1
                        stats['failed_file_paths'].append(_relative_name(file_path))

        stats['translation_stats'] = self.translation_service.get_stats()
        stats['vision_stats']      = self.vision_service.get_stats()
        stats['ocr_stats']         = self.ocr_service.get_stats()
        stats['total_estimated_cost'] = self.total_estimated_cost
        self._save_conversion_state()
        return stats


def run_conversion(input_dir: str, output_dir: str,
                   progress_callback=None) -> Dict[str, Any]:
    """
    Programmatic entry point used by the GUI and CLI.

    Args:
        input_dir: Path to the folder containing audit input files.
        output_dir: Path where converted .txt files will be written.
        progress_callback: Optional callable(message: str) for progress updates.

    Returns:
        stats dict from FileConverter.convert_all_files()
    """
    if not os.path.exists(input_dir):
        raise FileNotFoundError(f"Input directory not found: {input_dir}")

    os.makedirs(output_dir, exist_ok=True)

    # Redirect logger messages to the callback if provided
    cb_handler = None
    if progress_callback and not logging.getLogger().handlers:
        class _CBHandler(logging.Handler):
            def emit(self, record):
                progress_callback(self.format(record))
        cb_handler = _CBHandler()
        cb_handler.setFormatter(logging.Formatter('%(levelname)s: %(message)s'))
        logger.addHandler(cb_handler)

    try:
        converter = FileConverter(input_dir, output_dir)
        stats = converter.convert_all_files()
    finally:
        if cb_handler is not None:
            logger.removeHandler(cb_handler)

    return stats


def main():
    """CLI entry point — accepts optional positional args: input_dir output_dir"""
    import argparse
    logging.basicConfig(
        level=logging.INFO,
        format='%(asctime)s - %(levelname)s - %(message)s'
    )
    parser = argparse.ArgumentParser(description="Audit Tool — File Converter")
    parser.add_argument("input_dir", nargs="?", default=None,
                        help="Input folder with audit materials")
    parser.add_argument("output_dir", nargs="?", default=None,
                        help="Output folder for converted text files")
    args = parser.parse_args()

    input_dir = args.input_dir or os.environ.get("AUDIT_INPUT_DIR") or "test_input"
    output_dir = args.output_dir or os.environ.get("AUDIT_OUTPUT_DIR") or "converted_files"

    try:
        stats = run_conversion(input_dir, output_dir)
    except FileNotFoundError as e:
        logger.error(str(e))
        sys.exit(1)

    logger.info("\n=== CONVERSION STATISTICS ===")
    logger.info(f"Total files: {stats['total_files']}")
    logger.info(f"Converted: {stats['converted_files']}")
    logger.info(f"Failed: {stats['failed_files']}")
    logger.info(f"Skipped: {stats['skipped_files']}")
    logger.info(f"Unchanged: {stats.get('unchanged_files', 0)}")
    logger.info(f"Converted with warnings: {stats.get('warning_files', 0)}")
    if stats.get('unsupported_file_paths'):
        logger.warning("\n=== UNSUPPORTED FILES NOT PROCESSED ===")
        for file_path in stats['unsupported_file_paths']:
            logger.warning(f"- {file_path}")
    if stats.get('failed_file_paths'):
        logger.error("\n=== FILES NOT PROCESSED DUE TO ERRORS ===")
        for file_path in stats['failed_file_paths']:
            logger.error(f"- {file_path}")
    if stats.get('warning_file_paths'):
        logger.warning("\n=== FILES RETAINED WITH WARNINGS (WILL RETRY) ===")
        for file_path in stats['warning_file_paths']:
            logger.warning(f"- {file_path}")

    trans_stats = stats.get('translation_stats', {})
    logger.info("\n=== TRANSLATION STATISTICS ===")
    logger.info(f"Text checked: {trans_stats.get('total_text_checked', 0)}")
    logger.info(f"Text translated: {trans_stats.get('text_translated', 0)}")
    logger.info(f"Already English: {trans_stats.get('text_already_english', 0)}")
    logger.info(f"Translation errors: {trans_stats.get('translation_errors', 0)}")
    logger.info(f"Translation retries: {trans_stats.get('translation_retries', 0)}")
    logger.info(f"Translation skipped while unavailable: {trans_stats.get('translation_skipped_unavailable', 0)}")

    vision_stats = stats.get('vision_stats', {})
    provider = vision_stats.get('provider_used', 'none')
    logger.info("\n=== AI VISION STATISTICS ===")
    logger.info(f"Provider used: {provider.upper() if provider else 'None'}")
    logger.info(f"Images processed: {vision_stats.get('images_processed', 0)}")
    logger.info(f"API errors: {vision_stats.get('api_errors', 0)}")
    logger.info(f"Actual cost (approx): ${vision_stats.get('total_cost_actual', 0):.4f}")
    ocr_stats = stats.get('ocr_stats', {})
    logger.info("\n=== OCR STATISTICS ===")
    logger.info(f"OCR attempts: {ocr_stats.get('ocr_attempts', 0)}")
    logger.info(f"OCR errors: {ocr_stats.get('ocr_errors', 0)}")
    logger.info(f"OCR unavailable events: {ocr_stats.get('ocr_unavailable', 0)}")
    logger.info(f"\nTotal estimated cost: ${stats.get('total_estimated_cost', 0):.4f}")

    if stats['converted_files'] > 0 or stats.get('unchanged_files', 0) > 0:
        logger.info("\nConversion completed successfully!")
    else:
        logger.error("\nFailed to convert any files!")
        sys.exit(1)


if __name__ == "__main__":
    main()
